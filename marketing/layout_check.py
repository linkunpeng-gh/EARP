#!/usr/bin/env python3
"""layout_check.py — 检查 pptx 文本溢出与越界（几何启发式）
用法: python layout_check.py deck.pptx
"""
import sys, unicodedata
from pptx import Presentation
from pptx.util import Emu

EMU_PER_IN = 914400

def is_cjk(ch):
    o = ord(ch)
    return (0x4E00 <= o <= 0x9FFF) or (0x3000 <= o <= 0x303F) or (0xFF00 <= o <= 0xFFEF) \
        or (0x2000 <= o <= 0x206F and unicodedata.east_asian_width(ch) in ('W', 'F')) \
        or unicodedata.east_asian_width(ch) in ('W', 'F')

def char_w(ch, size_pt):
    """估计字符宽度（英寸）"""
    if ch in (' ', '\u3000'):
        return size_pt / 72.0 * (1.0 if ch == '\u3000' else 0.32)
    if is_cjk(ch):
        return size_pt / 72.0
    if ch.isdigit():
        return size_pt / 72.0 * 0.56
    if ch.isascii():
        return size_pt / 72.0 * 0.54
    return size_pt / 72.0 * 0.6

def wrap_lines(text, width_in, size_pt, bullet=False):
    usable = width_in - (0.22 if bullet else 0.0)
    lines = 0
    cur = 0.0
    for ch in text:
        w = char_w(ch, size_pt)
        if cur + w > usable and cur > 0:
            lines += 1
            cur = w
        else:
            cur += w
    if cur > 0 or text == '':
        lines += 1
    return lines

def para_height(par, width_in):
    txt = ''.join(r.text for r in par.runs)
    if not txt.strip():
        return 0.0
    size = None
    for r in par.runs:
        if r.font.size:
            size = r.font.size.pt
    if size is None:
        size = 14.0
    # lineSpacing
    ls = 1.0
    try:
        if par.line_spacing is not None:
            if isinstance(par.line_spacing, float):
                ls = par.line_spacing
            else:
                # spcPts 长度（EMU）→ 换算为倍数
                ls = (par.line_spacing / 12700.0) / size
    except Exception:
        pass
    bullet = par.level is not None and par.level >= 0 and 'buChar' in par._pPr.xml if par._pPr is not None else False
    lines = wrap_lines(txt, width_in, size, bullet)
    lh = size / 72.0 * 1.22 * max(ls, 0.8)
    return lines * lh

def check(path):
    prs = Presentation(path)
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            try:
                left, top, w, h = sh.left, sh.top, sh.width, sh.height
            except Exception:
                continue
            if left is None or w is None:
                continue
            L, T, W, H = left / EMU_PER_IN, top / EMU_PER_IN, w / EMU_PER_IN, h / EMU_PER_IN
            # 越界
            if L < -0.02 or T < -0.02 or L + W > 13.36 or T + H > 7.52:
                issues.append((si, 'OUT-OF-SLIDE', f'box=({L:.2f},{T:.2f},{W:.2f}x{H:.2f}) text={sh.text[:18]!r}'))
            # 溢出
            total_h = 0.0
            for p in sh.text_frame.paragraphs:
                total_h += para_height(p, W)
            if total_h > H + 0.06:
                issues.append((si, 'OVERFLOW', f'need={total_h:.2f} box={H:.2f} text={sh.text[:30]!r}'))
    return issues

if __name__ == '__main__':
    issues = check(sys.argv[1])
    if not issues:
        print('NO LAYOUT ISSUES FOUND')
    else:
        for si, kind, msg in issues:
            print(f'slide {si:02d} [{kind}] {msg}')
        print(f'TOTAL: {len(issues)}')
